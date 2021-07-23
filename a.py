import os
import shutil
import json
import xmltodict
import natsort
import zipfile
import re
import pathlib

from pptx import Presentation
from lxml import etree
from collections import OrderedDict


def unzip(src, des):
    """
    unzip the deck
    """
    with zipfile.ZipFile(src, 'r') as zip_ref:
        zip_ref.extractall(des)
    return


def zipdir(path, file_name):
    """
    zip extracted deck to get output deck
    """
    length = len(path)
    zipf = zipfile.ZipFile('output/'+f'Test_{file_name}.pptx', 'w', zipfile.ZIP_DEFLATED)
    for root, dirs, files in os.walk(path):
        folder = root[length:] # path without "parent"
        for file in files:
            zipf.write(os.path.join(root, file), os.path.join(folder, file))
    zipf.close()
    return


def gen_tree(path):
    """
    pass the path of the xml document to enable the parsing process
    """
    # print("CALLING.. Tree")
    parser = etree.XMLParser(remove_blank_text=True)
    tree = etree.parse(path, parser)
    root = tree.getroot()    
    return root, tree


def max_rId():
    """
    returns maximum rId
    """
    path = f'{output_path}/ppt/_rels/presentation.xml.rels'
    root, tree = gen_tree(path)
    
    rIds = []
    
    for relation in root:
        attrib = relation.attrib
        rId = int(attrib.get('Id').split('Id')[-1])
        rIds.append(rId)
    return {'rId': max(rIds)}


def xml_to_dict(path):
    """
    convert xml to dict
    """
    with open(path) as xml_file:
        data_dict = xmltodict.parse(xml_file.read())
        xml_file.close()
    if isinstance(data_dict["Relationships"]["Relationship"], list):
        data = sorted(data_dict["Relationships"]["Relationship"], key=lambda item: int(item['@Id'].split('Id')[1]))
    else:
        data = [data_dict["Relationships"]["Relationship"]]
    return data


def total_slides(path):
    """
    returns total number of slides
    """
    # print("CALLING.. total_slides")
    prs = Presentation(path)
    tot_slides = len(prs.slides._sldIdLst)
    return tot_slides


def first_slide(path):
    """
    returns the first slide rId
    """
    # print("CALLING.. First_slide")
    root, _ = gen_tree(path)

    slide = 'slide1.xml'
    for relation in root:
        attrib = relation.attrib
        if slide in attrib['Target']:
            return int(attrib['Id'].split('Id')[-1])


def get_rels(dict_1):
    """
    list latest .rels files
    """
    # print("DDD: ", dict_1)
    # a = dict_1.keys()
    a = [i for i in dict_1.values()]
    lis = natsort.natsorted([i for i in a if '_rels' in i])
    # print("HHH: ", lis)
    return lis


def add_slides(file_name):
    """
    returns a list of all slides
    """
    path = f'{input_decks}/{file_name}.pptx'
    sld = total_slides(path)+1
    lis = list(range(1, sld))
    return lis


def repair_path(dict_1):
    """
    replacing '\\' with '/'
    """
    dup_dict = dict(dict_1)
    for k,v in dup_dict.items():
        if '\\' in k:
            key = k.replace('\\', '/')
            val = v.replace('\\', '/')
            del dict_1[k]
            dict_1[key] = val
    return dict_1


def remove_dup(files1, dict_3, files2):
    """
    remove duplicates from the rels files
    """
    l1 = files1[:]
    for i in l1:
        if '/' not in i:
            if i in files2:
                files1.remove(i)
                del dict_3[i]
    return files1   


def new(path):
    """
    create, move and unzip the empty output deck
    """
    fq_empty = "resources/Empty.pptx"
    # create
    prs = Presentation()
    prs.save(fq_empty)
    # move
    d = ".".join ([path, "pptx"]) # "output/41.pptx"
    shutil.move (fq_empty, d)
    # unzip
    unzip (d, d.split('.')[0])
    os.remove(d)
    m_rId = max_rId()
    return m_rId


def make_dir(file_name): # output_file_loc = des
    """
    creates input deck directories which does not exists in the output deck
    """
    for i in os.walk(f'{tmp_path}/{file_name}'):
        fld = i[0].split(file_name)[-1]
        if fld:
            loc = f"{output_path}{fld}"
            if not os.path.exists(f'{output_path}/{fld}'):
                os.makedirs(f'{output_path}/{fld}')
    # print("MAKE_DIR completed...")        
    return


def make_structure(file_name):
    """
    creates structure of input deck
    """
    for i in os.walk(f'{tmp_path}/{file_name}'):
        fld = i[0].split(file_name)[-1]
        if fld:
            loc = f"{output_path}{fld}"
            if 'ppt' not in loc and (file_name not in loc):
                shutil.rmtree(f'{output_path}/{fld}')
                shutil.copytree(f'{tmp_path}/{file_name}/{i[0].split(file_name)[-1]}', f'{output_path}/{fld}')
    return


def add_files(path, file_name, target_files, slides=None):
    """
    returns a list of files that needs to be modified in output deck
    """
    # print("CALLING.. Add_files")
    data = xml_to_dict(path)
    
    if slides:
        sldIds = []
        
        tot_slides = total_slides(f'{input_decks}/{file_name}.pptx')
        first_slide_id = first_slide(path)
    
        files = []
        for i in data:
            current_rId = int(i['@Id'].split('Id')[1])
            if (first_slide_id > current_rId) or (current_rId > (first_slide_id+tot_slides-1)):
                if 'slideLayouts' not in i['@Target'] and 'slideMasters' not in i['@Target'] and 'theme' not in i['@Target']:
                    files.append(i['@Target'])
        for i in files:
            if '/' in i:
                a = i.split('/')
                fld, fl = a[0], a[1]
                if os.path.exists(f'{tmp_path}/{file_name}/ppt/{fld}/_rels'):
                    if os.path.isfile(f'{tmp_path}/{file_name}/ppt/{fld}/_rels/{fl}.rels'):
                        target_files.append(f'{fld}/_rels/{fl}.rels')
        
        target_files = target_files + files
        
        for id in slides:
            slide = f'slide{str(id)}.xml'
            sldIds.append([i["@Id"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            target_files.append([i["@Target"] for i in data if slide in i["@Target"] and "http" not in i["@Target"]][0])
            target_files.append(f'slides/_rels/{slide}.rels')
            add_files(f'{tmp_path}/{file_name}/ppt/slides/_rels/{slide}.rels', file_name, target_files)
    else:
        for i in data:
            # handling duplicacy
            if i["@Target"] in target_files or i["@Target"][3:] in target_files:
                pass
            elif "http" not in i["@Target"]:
                # if 'slideLayouts' not in i['@Target'] and 'slideMasters' not in i['@Target'] and 'theme' not in i['@Target']:
                if '../' in i["@Target"]:
                    target_files.append(i['@Target'][3:])
                else:
                    new_tar = path.split('/')[-3]
                    target_files.append(f'{new_tar}/{i["@Target"]}')
                
                if ".." in i['@Target'] and "xml" in i['@Target']:
                    path = f"{tmp_path}/{file_name}/ppt/{i['@Target'].split('..')[1].split('/')[1]}/_rels/{i['@Target'].split('..')[1].split('/')[2]}.rels"
                    
                    if os.path.exists(path):
                        # handling rels files
                        target_files.append(path.split('ppt/')[1])
                        add_files(path, file_name, target_files)

    return target_files


def get_fld_fl(file):
    """
    returns folder anme and file name
    """
    if '_rels' in file: # slides/_rels/slide2.xml.rels
        sp = file.split('/')
        fl_name = sp[-1]
        fld_name = f'{sp[0]}/{sp[1]}'
    elif '../' in file:
        _,fld_name,fl_name = file.split('/')
    else:
        fld_name,fl_name = file.split('/')
    
    return fld_name, fl_name


def list_target(target_files, d2):
    """
    creates a dict with number of files
    """
    count = 0
    for file in target_files:
        if '/' in file:
        # get folder and file name
            fld, fl = get_fld_fl(file)
            if fld not in d2:
                d2[fld] = 0
        if 'slideMasters' not in d2:
            d2['slideMasters'] = 0
        if 'slideLayouts' not in d2:
            d2['slideLayouts'] = 0
        if 'theme' not in d2:
            d2['theme'] = 0
    return d2
    

def rename(path, fld, fl, dict_2): # fld=media, fl=image1.png
    """
    rename a file
    """
    d1 = OrderedDict()
    
    ext = ''.join(pathlib.Path(fl).suffixes)
    name = re.findall(r'(\w+?)(\d+)', fl)[0][0]
    if f'{fld}/{name}' in dict_2.keys():
        count = dict_2[f'{fld}/{name}']+1
    else:
        count = 1
    new_name = f'{name}{count}{ext}'
    if 'ppt' in path:
        shutil.copy(f'{path}/{fld}/{fl}', f"{output_path}/ppt/{fld}/{new_name}")
    else:
        shutil.copy(f'{path}/{fld}/{fl}', f"{output_path}/{fld}/{new_name}")
    d1[f'{fld}/{fl}'] = f'{fld}/{new_name}'
    dict_2[f'{fld}/{name}'] = count
    return d1


def del_files(rels_fl, last_fl, path):
    """
    delete extra files
    Not using this funct anywhere
    """
    for i in rels_fl:
        if i[:-5] not in last_fl:
            os.remove(f'{path}/{i}')
    return


def copy_mandatory(src, des, deck, dict_1):
    """
    copy mandatory files
    """
    # print("COPY MANDATORY CALLING")
    m_list = ['slideLayouts', 'theme', 'slideMasters']
    if deck == 1:
        for fl in m_list:
            count = 0
            for i in os.walk(f'{src}/ppt/{fl}'):
                count = len(i[2])
            if os.path.exists(f'{des}/{fl}'):
                shutil.rmtree(f'{des}/{fl}')
            shutil.copytree(f'{src}/ppt/{fl}', f'{des}/{fl}')
            dict_2[fl] = count
            if os.path.exists(f'{src}/ppt/{fl}/_rels'):
                dict_2[f'{fl}/_rels'] = count
            
            for i in os.walk(f'{src}/ppt/{fl}'):
                fld = i[0].split('ppt/')[1]
                if '\\' in fld:
                    fld = fld.replace('\\', '/')
                fl_list = natsort.natsorted(i[2])
                for j in fl_list:
                    dict_1[f'{fld}/{j}'] = f'{fld}/{j}'
    else:
        for i in m_list:
            if os.path.exists(f'{des}/{i}'):
                for j in os.walk(f'{src}/ppt/{i}'):
                    fld = j[0].split('ppt/')[1]
                    if '\\' in fld:
                        fld = fld.replace('\\', '/')
                    
                    fl_list = natsort.natsorted(j[2])
                    for x in fl_list:
                        ext = ''.join(pathlib.Path(x).suffixes)
                        name = re.findall(r'(\w+?)(\d+)', x)[0][0]
                        count = dict_2[fld]+1
                        new_name = f'{name}{count}{ext}'
                        shutil.copy(f'{src}/ppt/{fld}/{x}', f'{des}/{fld}/{new_name}')
                        
                        dict_1[f'{i}/{x}'] = f'{fld}/{new_name}'
                        dict_2[fld] = count
    
    # remove empty folders
    for i in os.walk(des):
        if not i[2]:
            shutil.rmtree(i[0])
    
    return dict_1, dict_2


def copy_target(target_files, file_name, tmp_loc, dict_2):
    """
    copy target files from tmp to output folder 
    """
    d_1 = OrderedDict()
    
    target_files = natsort.natsorted(target_files)
    for i in target_files:
        if '/' in i:
            if 'slideLayouts' not in i and 'slideMasters' not in i and 'theme' not in i:
                fld_name,fl_name = get_fld_fl(i)
                if os.path.exists(f'{tmp_loc}/ppt/{fld_name}/{fl_name}'):
                    path = f'{tmp_loc}/ppt'
                    d_1.update(rename(path, fld_name, fl_name, dict_2))
                else:
                    d_1.update(rename(tmp_loc, fld_name, fl_name, dict_2))
        # else:
        #     if not os.path.isfile(f'{output_path}/ppt/{i}'):
        #         if os.path.isfile(f'{tmp_loc}/ppt/{i}'):
        #             shutil.copyfile(f'{tmp_loc}/ppt/{i}', f'{output_path}/ppt/{i}')
    # print("CALLING... copy_target")
    return d_1


def create_json(fl, name):
    """
    creates a json files
    """
    obj = json.dumps(fl, indent=4)
    with open(f"new_json/{name}.json", "w") as outfile:
        outfile.write(obj)
    return


def update_rels(fl_list, tmp_loc, dict_1):
    """
    update latest .rels files
    refactor the names of the assests and update the content
    """
    old_files = natsort.natsorted([i for i in dict_1.keys()])
    path = f'{output_path}/ppt'
    for i in fl_list:
        root, tree = gen_tree(f'{path}/{i}')
        for relation in root:
            attrib = relation.attrib
            if attrib.get('Target')[3:] in old_files:
                relation.set('Target', f"../{dict_1[attrib.get('Target')[3:]]}")
        tree.write(f'{path}/{i}', pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return


def get_relations(inp_path, file_name, slides):
    """
    returns f1:list of inp targets, d3:dict of targets, 
    f2: list of out target and sldIds: list of rIds of slides
    """
    root1,_ = gen_tree(inp_path)
    root2,_ = gen_tree(f'{output_path}/ppt/_rels/presentation.xml.rels')
    data = xml_to_dict(inp_path)
    tot_slides = total_slides(f'{input_decks}/{file_name}.pptx')
    first_slide_id = first_slide(inp_path)

    dict_3 = OrderedDict()
    files1 = []
    files2 = []
    sldIds = []

    for relation in root1:
        attrib = relation.attrib
        current_rId = int(attrib.get('Id').split('Id')[-1])
        if (first_slide_id > current_rId) or (current_rId > (first_slide_id+tot_slides-1)):
            files1.append(attrib["Target"])
            dict_3[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]
        if not slides:
            if (first_slide_id <= current_rId) and (current_rId < (first_slide_id+tot_slides)):
                files1.append(attrib["Target"])
                sldIds.append(attrib['Id'])
                dict_3[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]
    
    if slides:
        for id in slides:
            slide = f'slide{str(id)}.xml'
            for relation in root1:
                attrib = relation.attrib
                if slide in attrib['Target'] and "http" not in attrib['Target']:
                    files1.append(attrib['Target'])
                    sldIds.append(attrib['Id'])
                    dict_3[attrib['Target']] = [relation.tag, attrib['Id'], attrib['Type'], attrib['Target']]
    files1 = natsort.natsorted(files1)
    
    for relation in root2:
        attrib = relation.attrib
        files2.append(attrib['Target'])
    files2 = natsort.natsorted(files2)
    
    return files1, dict_3, files2, sldIds


def update_dict_3(dict_1, dict_3):
    """
    update dict_3 by removing '../' from the target
    """
    inp_keys = [i for i in dict_1.keys()]
    d3_keys = [i for i in dict_3.keys()]
    out_keys = natsort.natsorted([i for i in d3_keys])
    
    for i in out_keys:
        if '/' in i:
            val = dict_3[i]
            if '../' in i:
                val[3] = f'../{dict_1[i[3:]]}'
            else:
                val[3] = dict_1[i]
            dict_3[i] = val
    return dict_3


def update_rId(dict_2, files1, dict_3):
    """
    update the rIds and returns dict_2: modified max rId
    dict_3: updated rId for assests
    d1 = mapping of old and new rIds
    """
    d1 = OrderedDict()
    max_rId = dict_2['rId']
    for i in files1:
        max_rId += 1
        val = dict_3[i]
        d1[val[1]] = f'rId{max_rId}'
        val[1] = f'rId{max_rId}'
        dict_3[i] = val
        
    dict_2['rId'] = max_rId
    return dict_2, dict_3, d1


def write_rels(dict_3, files1):
    """
    adding assests in presentation.xml.rels
    by inserting assest relationship elements
    """
    path = f'{output_path}/ppt/_rels/presentation.xml.rels'
    root, tree = gen_tree(path)
    for i in files1:
        val = dict_3[i]
        tag, Id, Type, target = val
        ele = etree.Element(tag)
        etree.SubElement(root, tag, Id=Id, Type=Type, Target=target)
    tree.write(path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return
        

def xml_tag(inp_tag, out_tag):
    """
    returns a dict, key=new_tag, value=prev_tag
    """
    tag_dict = OrderedDict()
    sub_tag = OrderedDict()
    for i in range(len(inp_tag)):
        if inp_tag[i] not in out_tag:
            tag_dict[inp_tag[i]] = [inp_tag[i-1]]
    return tag_dict


def create_tags(tag_dict, o_tree):
    """
    creates empty tags in presentation.xml
    """
    for i, o in tag_dict.items():
        subtag1 = o_tree.find(o[0])
        subtag2 = etree.Element(i)
        subtag1.addnext(subtag2)
    o_tree.write(f'{output_path}/ppt/presentation.xml', pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    return 
   

def modify_d2(d1, d2):
    """
    modify the subtag dictionary
    """
    val_list = [i for i in d2.keys()]
    
    for key in val_list:
        for i in range(len(d2[key])):
            try:
                val = d1[d2[key][i][2]]
                d2[key][i][2] = val
                if None in d2[key][i]:
                    d2[key][i].remove(None)
            except:
                pass
    return d2
            
            
def add_extLst(src_xml, des_xml, ext_lst, tag_dict):
    """
    adding extlst subelements in prep.xml file
    """
    inp_root,_ = gen_tree(src_xml)
    out_root, out_tree = gen_tree(des_xml)
    
    for relation in ext_lst:
        
        # if relation in tag_dict.keys():
        #     print("JJJ: ", relation)
        #     print("PPP: ", tag_dict[relation])
        for elt in inp_root.findall(relation):
            # print("ELE: ", elt.tag)
            out_root.append(elt)

    out_tree.write(des_xml, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return
         

def get_prep_tags(src_xml, d1):
    """
    get tags and subtags from input presentation.xml
    returns a dict with tags and subtags for presentation.xml
    """
    root, tree = gen_tree(src_xml)
    d2 = OrderedDict()
    # list of old rIds
    rId_lis = [i for i in d1.keys()]
    nmsps =  root.nsmap['r']
    ext_lst = []
    
    for relation in root:
        for ele in relation:
            attrib = ele.attrib
            tag = ele.tag
            # try:
            # if attrib[f"{{{nmsps}}}id"]:
            if attrib.get(f"{{{nmsps}}}id"):
                if attrib.get(f"{{{nmsps}}}id") in rId_lis:
                    if relation.tag in d2:
                        val = d2[relation.tag]
                        val.append([tag, attrib.get('id'), attrib.get(f"{{{nmsps}}}id")])
                        d2[relation.tag] = val
                    else:
                        d2[relation.tag] = [[tag, attrib.get('id'), attrib.get(f"{{{nmsps}}}id")]]
            else:
                if 'uri' in ele.attrib:
                    if relation.tag not in ext_lst:
                        ext_lst.append(relation.tag)
                    # print("ELE11: ", ele)
                    # extLst.append(ele)
                    # if relation.tag in d2:
                    #     val = d2[relation.tag]
                    #     val.append(ele)
                    #     d2[relation.tag] = val
                    # else:
                    #     d2[relation.tag] = [ele]
    d2 = modify_d2(d1, d2)
    return d2, ext_lst


def add_subtags(path, pxml_subtags):
    """
    add subtags in the presentation.xml file
    """
    # print("CALLING... add_subtag")
    root, tree = gen_tree(path)
    nmsps =  root.nsmap['r']
    for k,v in pxml_subtags.items():
        subtag1 = tree.find(k)
        for i in v:
            if 'rId' not in i[1]:
                rId = f"{{{nmsps}}}id"
                subtext = etree.SubElement(subtag1, i[0])
                subtext.attrib['id'] = i[1]
                subtext.attrib[rId] = i[2]
            else:
                subtext = etree.SubElement(subtag1, i[0])
                subtext.attrib[rId] = i[1]
    tree.write(f'{output_path}/ppt/presentation.xml', pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return


def clean_prep_xml(des_xml, rels_rIds, pxml_subtags):
    """
    remove extra subtags if their rId is not present in 
    the presentation.xml.rels file
    """
    root, tree = gen_tree(des_xml)
    nmsps =  root.nsmap['r']
    rId = f"{{{nmsps}}}id"
    for k,v in pxml_subtags.items():
        subtag1 = tree.find(k)
        for i in subtag1:
            if i.attrib.get(rId):
                if i.attrib.get(rId) not in rels_rIds:
                    subtag1.remove(i)

    tree.write(des_xml, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return


def rel_duplicates():
    """
    remove the duplicates entries of 'Target' field 
    form presentation.xml.rels if any and 
    returns a dict of key:target, val:rId 
    """
    path = f'{output_path}/ppt/_rels/presentation.xml.rels'
    root, tree = gen_tree(path)
    d1 = OrderedDict()
    for relation in root:
        rIds = []
        attrib = relation.attrib
        if attrib['Target'] in d1.keys():
            val = d1[attrib['Target']]
            val.append(attrib['Id'])
            d1[attrib['Target']] = val
        else:
            d1[attrib['Target']] = [attrib['Id']]
    
    # getting duplicates rIds
    dup_rIds = []
    for k,v in d1.items():
        if len(v) > 1:
            dup_rIds.append(v.pop(0))
            d1[k] = v
    
    # removing relation
    for relation in root:
        attrib = relation.attrib
        if attrib['Id'] in dup_rIds:
            root.remove(relation)
    
    rels_rIds = [relation.attrib['Id'] for relation in root]
        
    tree.write(path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return d1, rels_rIds


def scan_sldsz(src_xml, des_xml):
    """
    scan slide size for templating
    """
    _,i_tree = gen_tree(src_xml)
    _,o_tree = gen_tree(des_xml)
    
    tag = "{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz"
    
    inp_sldsz = i_tree.find(tag).attrib
    cx = inp_sldsz['cx']
    cy = inp_sldsz['cy']
    Type = inp_sldsz.get('type')
    out_sldsz = o_tree.find(tag).attrib
    if out_sldsz['cx']!=cx or out_sldsz['cy']!=cy:
        out_sldsz['cx'] = cx
        out_sldsz['cy'] = cy
    if not Type and out_sldsz.get('type'): # error
        del out_sldsz['type']
    o_tree.write(des_xml, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return


def write_pres(tmp_loc, d1, rels_rIds):
    """
    update the presentation.xml file
    """
    src_xml = f'{tmp_loc}/ppt/presentation.xml'
    des_xml = f'{output_path}/ppt/presentation.xml'
    
    inp_root, inp_tree = gen_tree(src_xml)
    out_root, out_tree = gen_tree(des_xml)
    
    inp_tag = [relation.tag for relation in inp_root]
    out_tag = [relation.tag for relation in out_root]
    
    tag_dict = xml_tag(inp_tag, out_tag)
    create_tags(tag_dict, out_tree)
    
    pxml_subtags, ext_lst = get_prep_tags(src_xml, d1)
    # add_extLst(src_xml, des_xml, ext_lst, tag_dict)
    create_json(pxml_subtags, '05_pxml_subtag')
    clean_prep_xml(des_xml, rels_rIds, pxml_subtags)
    add_subtags(des_xml, pxml_subtags)
    scan_sldsz(src_xml, des_xml)
    return


def presenation_files(inp_pres_rels, file_name, slides, dict_1, dict_2, tmp_loc):
    """
    deals with rels and xml file of presentation
    """
    files1, dict_3, files2, sldIds = get_relations(inp_pres_rels, file_name, slides)
    files1 = remove_dup(files1, dict_3, files2)
    # print("LLL: ", files1)
    dict_3 = update_dict_3(dict_1, dict_3)
    dict_2, dict_3, d1 = update_rId(dict_2, files1, dict_3)
    create_json(dict_3, '03_prepRelSubtag')
    write_rels(dict_3, files1)
    prep_rels_rIds, rels_rIds = rel_duplicates()
    create_json(prep_rels_rIds, '04_prep_rels_rIds')
    
    write_pres(tmp_loc, d1, rels_rIds)
    return dict_2


def handle_configs(tmp_loc):
    """
    handle configuration files
    """
    inp_path = '/'.join([tmp_loc, 'ppt'])
    out_path = f'{output_path}/ppt'
    
    config_fls = [i for i in os.listdir(inp_path) if os.path.isfile(f'{inp_path}/{i}')]
    # print("CCC: ", config_fls)
    mergables = ['commentAuthors.xml', 'tableStyles.xml']
    sing_prop = ['viewProps.xml', 'presProps.xml']
    ignore = ['revisionInfo.xml']
    
    for i in config_fls:
        inp_fl = f'{inp_path}/{i}'
        out_fl = f'{out_path}/{i}'
        # if i in mergables:
        if os.path.isfile(f'{out_path}/{i}'):
            root1,tree1 = gen_tree(inp_fl)
            root2,tree2 = gen_tree(out_fl)
            if i in mergables:
                try:
                    for relation in [f"{root1[0].tag}"]:
                        for elt in root1.findall(relation):
                            root2.append(elt)
                except:
                    pass
            elif i in sing_prop:
                if i == 'presProps.xml':
                    inp_d = {}
                    out_lis = []
                    nm = root1.nsmap['p']
                    tag0 = f"{{{nm}}}extLst"
                    for relation in [f"{root1[0].tag}"]:
                        fp = root1.find(tag0)
                        for ele in fp:
                            attrib = ele.attrib
                            if attrib['uri'] not in inp_d.keys():
                                inp_d[attrib['uri']] = ele
                    
                    for relation in [f"{root2[0].tag}"]:
                        fp = root2.find(tag0)
                        for ele in fp:
                            attrib = ele.attrib
                            out_lis.append(attrib['uri'])
                    for k,v in inp_d.items():
                        if k not in out_lis:
                            tag1 = root2.find(tag0)
                            tag1.append(v)
            else:
                pass
            tree2.write(out_fl, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
        else:
            shutil.copyfile(inp_fl, out_fl)
        
            
    return


def handle_cleaning():
    """
    remove unnecessary files 
    like changesInfos, printerSettings
    """
    extra_fl = ['changesInfos', 'printerSettings']
    fld_path = f'{output_path}/ppt'
    out_rel_path = f'{fld_path}/_rels/presentation.xml.rels'
    root, tree = gen_tree(out_rel_path)
    
    for i in extra_fl:
        path = f'{fld_path}/{i}'
        if os.path.isfile(path):
            shutil.rmtree(path)
    
        for relation in root:
            attrib = relation.attrib
            if i in attrib['Target']:
                root.remove(relation)
    
    tree.write(out_rel_path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    return
    

def content(tmp_loc, ref_names_dict, order):
    """
    add content_type in [Contant_Types].xml file
    """
    
    fl = '[Content_Types].xml'
    inp_path = '/'.join([tmp_loc, fl])
    out_path = '/'.join([output_path, fl])
    
    cnt_lst = []
    asset_lst = []
    def_att = []
    d = dict()
    
    root1,tree1 = gen_tree(inp_path)
    root2,tree2 = gen_tree(out_path)
    
    # get all the extensions belongs to "Default" tag
    for relation in root2:
        if 'Default' in relation.tag:
            def_att.append(relation.attrib['Extension'])
        else:
            break
    
    for relation in root1:
        if 'Override' in relation.tag:
            attrib = relation.attrib['PartName'][1:]
            try:
                cnt = attrib.split('ppt/')[-1]
                ini = '/ppt/'
            except:
                cnt = attrib
                ini = '/'
            if cnt in ref_names_dict.keys():
                relation.attrib['PartName'] = f'{ini}{ref_names_dict[cnt]}'
                cnt_lst.append(relation)
                # asset_lst.append(relation.attrib['PartName'])
            else:
                cnt_lst.append(relation)
            if relation.attrib['PartName'] not in asset_lst:
                asset_lst.append(relation.attrib['PartName'])
        else:
            attrib = relation.attrib['Extension']
            if attrib not in def_att:
                cnt_lst.append(relation)
                # asset_lst.append(relation.attrib['Extension'])
        # deal with the assest_lst
    # print("AA: ", asset_lst)
    cnt_lst = natsort.natsorted(cnt_lst)
    for ele in cnt_lst:
        prev = tree2.find(ele.tag)
        prev.addnext(ele)
    
    tree2.write(out_path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)
    
    unq_attr = []
    for relation in root2:
        if 'Override' in relation.tag:
            if relation.attrib['PartName'] not in unq_attr:
                unq_attr.append(relation.attrib['PartName'])
            else:
                root2.remove(relation)
    tree2.write(out_path, pretty_print=True, xml_declaration=True, encoding='UTF-8', standalone=True)


def get_order(tmp_loc, order):
    """
    scans the presentation.xml.rels
    returns the order of the assets
    """
    fn = 'ppt/_rels/presentation.xml.rels'
    inp_path = f'{tmp_loc}/{fn}'
    out_path = f'{output_path}/{fn}'
    
    root1,tree1 = gen_tree(inp_path)
    root2,tree2 = gen_tree(out_path)
    
    rid_dict = {}
    for relation in root1:
        Type = relation.attrib['Type'].split('/')[-1]
        rId = int(relation.attrib['Id'].split('Id')[-1])
        rid_dict[rId] = Type
    
    ord_dict = OrderedDict(sorted(rid_dict.items()))
    type_dict = OrderedDict()
    prev = ''
    for k,v in ord_dict.items():
        if v not in type_dict.keys():
            type_dict[v] = [k]
        else:
            type_dict[v] = type_dict[v] + [k]
        if v not in order:
            if prev not in order:
                ind = 0
            else:
                ind = order.index(prev)
            order.insert(ind+1, v)
        prev = v
    # print("OOO: ", order)
    print("ORDER: ", order)
    return order, type_dict
        
    
def refactoring_rIds(inp_path, order, type_dict, dict_2):
    """
    refactors rIds of the assets
    """
    fl_lis = []
    for i in fl_lis:
        if i==order[0]:
            pass
    return
            

def deck_handler(id, msg, deck, dict_2, order):
    """
    handle the deck and select files for output deck
    """
    file_name, slides = msg['d'], msg['s']
    target_files = []
    
    tmp_loc = f'{tmp_path}/{file_name}'
    
    # unzip the input deck
    unzip(f'{input_decks}/{file_name}.pptx', tmp_loc)
    
    # creates folder structure of the input deck
    make_dir(file_name)
    inp_pres_rels = f'{tmp_loc}/ppt/_rels/presentation.xml.rels'
    
    if deck == 1:
        make_structure(file_name)
    
    # get all the slides if slide is None
    if not slides:
        slides = add_slides(file_name)
    
    target_files = add_files(inp_pres_rels, file_name, target_files, slides)
    # print("TARGET: ", target_files)
    # dict_2.update(list_target(target_files, dict_2))
    
    order, a = get_order(tmp_loc, order)
    # refactoring_rIds(inp_pres_rels)
    
    # copy all the required assets
    dict_1 = copy_target(target_files, file_name, tmp_loc, dict_2)
    # copy all the mendatory assets
    d1, d2 = copy_mandatory(tmp_loc, f'{output_path}/ppt', deck, dict_1)
    # updating the dictionaries
    dict_1.update(d1)
    dict_2.update(d2)

    # modify the rels files
    # get the latest rels file list
    rels_list = get_rels(dict_1)
    # refactor names in rels files
    update_rels(rels_list, tmp_loc, dict_1)
    # updating rels and xml files of presentation
    dict_2 = presenation_files(inp_pres_rels, file_name, slides, dict_1, dict_2, tmp_loc)
    # adding contents in the Content-Type file
    content(tmp_loc, dict_1, order)
    # handling the properties files
    handle_configs(tmp_loc)
    # removing extra files
    handle_cleaning()
    
    # creating json files
    create_json(dict_1, '01_refactored_names')
    create_json(dict_2, '02_refactoring_count')
    
    return order


if __name__ == '__main__':
    
    base_path = os.path.dirname(os.path.realpath(__file__))
    print("CURRENT_DIR:", base_path)
    dict_1 = OrderedDict()
    dict_2 = OrderedDict()
    order = []
    
    sample_msg = [41, {'d': 'Onboarding', 's':  [2,4,6]}, {'d': 'Presentation1','s':  None}]
    # sample_msg = [41, {'d': 'Onboarding', 's':  [2,4,6]}, {'d': 'Presentation1','s':  [1]}]
    # sample_msg = [41, {'d': 'Onboarding', 's':  [2, 4, 6]}]
    # sample_msg = [41, {'d': 'Onboarding', 's':  None}]
    # sample_msg = [41, {'d': 'Presentation1', 's': None}]
    # sample_msg = [41, {'d': 'Presentation1', 's':  [1]}] # working
    # sample_msg = [41, {'d': 'BI Case Studies', 's':  [2, 3]}]

    render_id = sample_msg.pop(0)
    
    output_path = f'{base_path}/output/{str(render_id)}'
    tmp_path = f'{base_path}/tmp/{render_id}'
    input_decks = f'{base_path}/presentations' 

    try:
        os.makedirs(output_path)
        os.makedirs(tmp_path)
    except:
        print("DIR ALREADY EXIST")
    
    # creating and unzipping empty deck
    m_rId = new(output_path)
    dict_2.update(m_rId)
    
    # iterating all the messages
    deck = 1
    while sample_msg:
        order = deck_handler(render_id, sample_msg.pop(0), deck, dict_2, order)
        deck += 1

    # zip the output folder
    zipdir(f'{output_path}', "Test")
